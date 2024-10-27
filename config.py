
import asyncio
import logging
from pathlib import Path
from typing import List, Optional

from aiogram import Bot, Dispatcher, types
from aiogram.filters import Command
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, Message, CallbackQuery
from pptx import Presentation
from pptx.util import Inches, Pt
import aiohttp
import aiofiles
import os
from PIL import Image
import io

# Настройка логирования
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Конфигурация
class Config:
    API_TOKEN = os.getenv("API_TOKEN", "7567293644:AAHVSPgyYAPt_NaUahhdID2njMK-FrtxaRg")  # Чтение токена из переменной окружения для безопасности
    MAX_FILE_SIZE = 500 * 1024 * 1024  # 500 MB
    DOWNLOAD_PATH = Path("downloads")
    TEMPLATES_PATH = Path("templates")

    # Создаем необходимые директории при запуске
    DOWNLOAD_PATH.mkdir(parents=True, exist_ok=True)
    TEMPLATES_PATH.mkdir(parents=True, exist_ok=True)


class PresentationManager:
    """Класс для управления презентациями"""

    @staticmethod
    async def create_presentation(template_name: Optional[str] = None) -> Presentation:
        """Создает новую презентацию или загружает шаблон"""
        if template_name:
            template_path = Config.TEMPLATES_PATH / template_name
            if template_path.exists():
                return Presentation(template_path)
            else:
                raise FileNotFoundError(f"Шаблон {template_name} не найден.")
        return Presentation()

    @staticmethod
    async def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
        """Добавляет титульный слайд"""
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if subtitle:
            title_slide.placeholders[1].text = subtitle

    @staticmethod
    async def add_content_slide(prs: Presentation, title: str, content: List[str]) -> None:
        """Добавляет слайд с контентом"""
        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        bullet_slide.shapes.title.text = title

        tf = bullet_slide.shapes.placeholders[1].text_frame
        for point in content:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)

    @staticmethod
    async def add_image_slide(prs: Presentation, title: str, image_path: str) -> None:
        """Добавляет слайд с изображением"""
        img_slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_slide.shapes.title.text = title

        # Добавляем изображение
        img_path = Path(image_path)
        if img_path.exists():
            left = top = Inches(1)
            img_slide.shapes.add_picture(str(img_path), left, top)
        else:
            raise FileNotFoundError(f"Изображение по пути {image_path} не найдено.")


class PresentationBot:
    """Основной класс бота для работы с презентациями"""

    def __init__(self):
        self.bot = Bot(token=Config.API_TOKEN)
        self.dp = Dispatcher()
        self.presentation_manager = PresentationManager()
        self.register_handlers()

    def register_handlers(self):
        """Регистрация обработчиков команд и нажатий на кнопки"""
        self.dp.message.register(self.cmd_start, Command("start"))
        self.dp.message.register(self.handle_document, lambda m: m.document is not None)
        # Регистрируем обработку нажатий на инлайн-кнопки
        self.dp.callback_query.register(self.handle_template_choice, lambda c: c.data.startswith('template_'))
        self.dp.callback_query.register(self.handle_callback_queries)  # Обработка всех остальных callback

    async def handle_callback_queries(self, callback: CallbackQuery):
        """Обработчик всех нажатий на инлайн-кнопки"""
        if callback.data == "upload_presentation":
            await callback.message.answer("Вы выбрали загрузку презентации. Пожалуйста, загрузите файл.")
        elif callback.data == "choose_template":
            await callback.message.answer("Выберите шаблон для вашей презентации.")
        elif callback.data == "create_new":
            await callback.message.answer("Создание новой презентации.")
        elif callback.data == "help":
            await callback.message.answer("Это бот для создания и редактирования презентаций.")
        else:
            await callback.message.answer("Неверная команда. Попробуйте снова.")

    # ... other methods ...

    @staticmethod
    def get_start_keyboard() -> InlineKeyboardMarkup:
        """Создает клавиатуру для начального меню"""
        keyboard = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📤 Загрузить презентацию", callback_data="upload_presentation")],
            [InlineKeyboardButton(text="🎨 Выбрать шаблон", callback_data="choose_template")],
            [InlineKeyboardButton(text="✍️ Создать с нуля", callback_data="create_new")],
            [InlineKeyboardButton(text="❓ Помощь", callback_data="help")]
        ])
        return keyboard

    async def cmd_start(self, message: Message):
        """Обработчик команды /start"""
        await message.answer(
            "👋 Добро пожаловать в Presentation Assistant Bot!\n\n"
            "Я помогу вам создать или отредактировать презентацию. Выберите действие:",
            reply_markup=self.get_start_keyboard()
        )

    async def handle_document(self, message: Message):
        """Обработка загруженной презентации"""
        document = message.document

        if not document:
            return

        # Логирование файла
        logger.info(f"Получен файл от {message.from_user.id}: {document.file_name} ({document.file_size} байт)")

        if document.file_size > Config.MAX_FILE_SIZE:
            await message.answer(f"⚠️ Файл слишком большой. Максимальный размер: {Config.MAX_FILE_SIZE // (1024 * 1024)} МБ")
            return

        try:
            file_id = document.file_id
            file = await self.bot.get_file(file_id)

            # Создаем уникальное имя файла
            file_path = Config.DOWNLOAD_PATH / f"{message.from_user.id}_{document.file_name}"

            # Скачиваем файл
            async with aiofiles.open(file_path, 'wb') as f:
                await self.bot.download_file(file.file_path, f)

            # Открываем презентацию для проверки
            prs = Presentation(file_path)

            await message.answer(
                f"✅ Презентация успешно загружена!\n"
                f"📊 Количество слайдов: {len(prs.slides)}\n\n"
                "Выберите действие:",
                reply_markup=self.get_edit_keyboard()
            )

        except Exception as e:
            logger.error(f"Ошибка при обработке презентации: {e}")
            await message.answer("❌ Произошла ошибка при обработке презентации. Попробуйте другой файл.")

    @staticmethod
    def get_edit_keyboard() -> InlineKeyboardMarkup:
        """Создает клавиатуру для редактирования презентации"""
        keyboard = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="➕ Добавить слайд", callback_data="add_slide")],
            [InlineKeyboardButton(text="🖼 Добавить изображение", callback_data="add_image")],
            [InlineKeyboardButton(text="📝 Изменить текст", callback_data="edit_text")],
            [InlineKeyboardButton(text="🎨 Изменить дизайн", callback_data="change_design")],
            [InlineKeyboardButton(text="💾 Сохранить", callback_data="save")]
        ])
        return keyboard

    async def handle_template_choice(self, callback: CallbackQuery):
        """Обработка выбора шаблона"""
        template_name = callback.data.replace('template_', '')
        try:
            prs = await self.presentation_manager.create_presentation(f"{template_name}.pptx")
            # Сохраняем временную презентацию
            temp_path = Config.DOWNLOAD_PATH / f"{callback.from_user.id}_temp.pptx"
            prs.save(temp_path)

            await callback.message.answer(
                "✅ Шаблон выбран успешно! Теперь вы можете:\n"
                "1. Добавить контент\n"
                "2. Изменить дизайн\n"
                "3. Сохранить презентацию",
                reply_markup=self.get_edit_keyboard()
            )
        except FileNotFoundError as e:
            logger.error(f"Ошибка создания презентации из шаблона: {e}")
            await callback.message.answer(f"❌ Ошибка: {str(e)}")
        except Exception as e:
            logger.error(f"Непредвиденная ошибка: {e}")
            await callback.message.answer("❌ Произошла ошибка при создании презентации")

    async def run(self):
        """Запуск бота"""
        logger.info("Запуск бота...")
        try:
            await self.dp.start_polling(self.bot)
        except Exception as e:
            logger.error(f"Ошибка при запуске бота: {e}")
        finally:
            await self.bot.session.close()


if __name__ == "__main__":
    bot = PresentationBot()
    asyncio.run(bot.run())

